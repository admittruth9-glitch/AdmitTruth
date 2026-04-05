# First, run: pip install python-pptx
from pptx import Presentation

prs = Presentation()

slides_content = [
    ("AdmitTruth – Website Project", "Transparent Admission Gateway\nFocus: Student-Centric Data Transparency"),
    ("Introduction", "• Centralized platform for reliable college data\n• Problem: Fragmented/biased information\n• Solution: User-friendly, accurate, and simple website"),
    ("Objectives", "• Authenticity: Verified information\n• Efficiency: Reduce time and confusion\n• Guidance: Help choose the right college\n• Usability: Accessible design"),
    ("Technologies Used", "• HTML5: Core structure\n• CSS3: Styling and responsive design\n• JavaScript: Interactivity\n• Responsive Design: Mobile and desktop optimization"),
    ("Homepage Overview", "• Clean Interface: Minimalist design\n• Navigation-Centric: Immediate access\n• Professional aesthetic builds trust"),
    ("Navigation Architecture", "• Links: Home, About, Colleges, Contact\n• Responsive Sidebar/Header\n• Logical user flow"),
    ("Key Features", "• Verified Data: Single source of truth\n• Mobile-First: Optimized for students\n• Performance: Fast loading speeds\n• Intuitive UI: Zero learning curve"),
    ("Functional Workflow", "1. Visit: Landing page arrival\n2. Browse: Explore databases\n3. Select: Filter based on needs\n4. Action: Gain accurate data"),
    ("Project Modules", "• Home: General overview\n• About: Mission and background\n• Colleges: Main database\n• Contact: Support channel"),
    ("Advantages", "• Time-Saving\n• Better Decisions\n• Reduced Stress\n• 24/7 Accessibility"),
    ("Current Limitations", "• Data Scale: Limited initial dataset\n• Connectivity: Requires internet\n• Interactivity: Static information in v1.0"),
    ("Future Roadmap", "• Expanded Database: More colleges\n• Search & Filters: Deep-search capabilities\n• AI Integration: Chatbot assistance\n• Personalization: User accounts"),
    ("Conclusion", "• Successfully developed AdmitTruth prototype\n• Met transparency objectives\n• Proven utility for streamlining admissions"),
    ("Thank You", "Questions & Discussion\n[Insert Contact Info]")
]

for title, content in slides_content:
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

prs.save('AdmitTruth_Presentation.pptx')